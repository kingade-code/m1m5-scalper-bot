from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG = RGBColor(0x0F, 0x17, 0x29)
GOLD = RGBColor(0xFF, 0xD7, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x00, 0xC8, 0x53)
RED = RGBColor(0xFF, 0x45, 0x45)
CYAN = RGBColor(0x00, 0xBF, 0xFF)
GRAY = RGBColor(0x8A, 0x93, 0xA6)
DARK = RGBColor(0x1A, 0x24, 0x3B)

def set_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG

def add_box(slide, left, top, w, h, fill_color=DARK):
    shp = slide.shapes.add_shape(1, left, top, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    shp.line.fill.background()
    return shp

def add_text(slide, left, top, w, h, text, sz=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align

def add_metric(slide, left, top, label, value, color=GREEN):
    add_box(slide, left, top, Inches(2.4), Inches(1.7), DARK)
    add_text(slide, left + Inches(0.1), top + Inches(0.15), Inches(2.2), Inches(0.4), label, 13, GRAY, False, PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.1), top + Inches(0.6), Inches(2.2), Inches(0.8), value, 30, color, True, PP_ALIGN.CENTER)

months_labels = ["Jun'26", "Jul", "Aug"]
equity_vals = [10495, 43150, 109423]
monthly_pnl = [10495, 32655, 66273]

# ── Slide 1: Title ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, Inches(0), Inches(1.5), Inches(13.333), Inches(1.5), "M1-M5 SCALPER BOT", 54, GOLD, True, PP_ALIGN.CENTER)
add_text(slide, Inches(0), Inches(3.0), Inches(13.333), Inches(1.0), "Backtest Report - Optimized Settings", 28, WHITE, False, PP_ALIGN.CENTER)
add_text(slide, Inches(0), Inches(4.0), Inches(13.333), Inches(0.8), "EMA 30 + Trailing Step 0.2  |  3-Month Performance", 20, GRAY, False, PP_ALIGN.CENTER)
add_text(slide, Inches(0), Inches(5.5), Inches(13.333), Inches(0.6), "XAUUSD  |  GBPUSD  |  AUDUSD  |  M1 + M5", 18, CYAN, False, PP_ALIGN.CENTER)
add_text(slide, Inches(0), Inches(6.5), Inches(13.333), Inches(0.5), "Jun 2026 - Aug 2026", 16, GRAY, False, PP_ALIGN.CENTER)

# ── Slide 2: Key Metrics ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, Inches(0), Inches(0.2), Inches(13.333), Inches(0.7), "KEY PERFORMANCE METRICS", 36, GOLD, True, PP_ALIGN.CENTER)
metrics = [
    ("WIN RATE", "66.1%", GREEN), ("PROFIT FACTOR", "2.36", GREEN),
    ("SHARPE RATIO", "5.13", GREEN), ("TOTAL TRADES", "10,529", CYAN),
    ("TOTAL P/L", "+$109,423", GREEN),
    ("MAX DRAWDOWN", "12.24%", GREEN), ("AVG WIN", "$27.29", GREEN),
    ("AVG LOSS", "-$22.52", RED), ("EXPECTANCY", "$10.39", GREEN),
    ("RECOVERY FACTOR", "261.72", GREEN),
]
for i, (label, value, color) in enumerate(metrics):
    row, col = divmod(i, 5)
    x = Inches(0.5) + Inches(col * 2.55)
    y = Inches(1.2) + Inches(row * 2.1)
    add_metric(slide, x, y, label, value, color)

# ── Slide 3: Equity Curve ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, Inches(0), Inches(0.2), Inches(13.333), Inches(0.7), "EQUITY CURVE", 36, GOLD, True, PP_ALIGN.CENTER)
fig, ax = plt.subplots(figsize=(11, 4.5))
fig.patch.set_facecolor('#0F1729')
ax.set_facecolor('#0F1729')
ax.fill_between(range(len(equity_vals)), equity_vals, alpha=0.3, color='#FFD700')
ax.plot(equity_vals, color='#FFD700', linewidth=2.5)
ax.set_xticks(range(len(months_labels)))
ax.set_xticklabels(months_labels, color='#8A93A6', fontsize=10)
ax.tick_params(axis='y', colors='#8A93A6')
ax.set_ylabel('Equity ($)', color='#8A93A6', fontsize=11)
ax.spines['bottom'].set_color('#8A93A6')
ax.spines['left'].set_color('#8A93A6')
ax.spines['top'].set_visible(False)
ax.spines['right'].set_visible(False)
ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, p: f'${x:,.0f}'))
plt.tight_layout()
plt.savefig('m1m5_equity.png', dpi=150, facecolor='#0F1729', bbox_inches='tight')
plt.close()
slide.shapes.add_picture('m1m5_equity.png', Inches(1.2), Inches(1.2), Inches(11), Inches(5.8))

# ── Slide 4: By Symbol ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, Inches(0), Inches(0.2), Inches(13.333), Inches(0.7), "PERFORMANCE BY SYMBOL", 36, GOLD, True, PP_ALIGN.CENTER)
symbols = ["XAUUSD", "GBPUSD", "AUDUSD"]
wr_sym = [65.3, 66.3, 66.4]
trades_sym = [2734, 3993, 3802]
fig, axes = plt.subplots(1, 3, figsize=(11, 4.5))
fig.patch.set_facecolor('#0F1729')
for i, (sym, wr, trades) in enumerate(zip(symbols, wr_sym, trades_sym)):
    ax = axes[i]
    ax.set_facecolor('#1A243B')
    w = trades * wr / 100
    l = trades * (1 - wr / 100)
    bars = ax.bar(['Wins', 'Losses'], [w, l], color=['#00C853', '#FF4545'], width=0.6)
    ax.set_title(f'{sym}\nWR: {wr}% | {trades} trades', color='white', fontsize=13, fontweight='bold', pad=10)
    ax.tick_params(colors='#8A93A6')
    for spine in ax.spines.values(): spine.set_color('#8A93A6')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, p: f'{x:,.0f}'))
    for bar in bars:
        ax.text(bar.get_x() + bar.get_width()/2., bar.get_height() + 30, f'{bar.get_height():,.0f}', ha='center', va='bottom', color='white', fontsize=11, fontweight='bold')
plt.tight_layout()
plt.savefig('m1m5_symbols.png', dpi=150, facecolor='#0F1729', bbox_inches='tight')
plt.close()
slide.shapes.add_picture('m1m5_symbols.png', Inches(1.2), Inches(1.2), Inches(11), Inches(5.8))

# ── Slide 5: By Timeframe ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, Inches(0), Inches(0.2), Inches(13.333), Inches(0.7), "PERFORMANCE BY TIMEFRAME", 36, GOLD, True, PP_ALIGN.CENTER)
fig, axes = plt.subplots(1, 2, figsize=(10, 4.5))
fig.patch.set_facecolor('#0F1729')
tfs = ["M1", "M5"]
wr_tf = [64.6, 68.0]
trades_tf = [5981, 4548]
pl_tf = [47059, 62364]
for i, (tf, wr, trades, pl) in enumerate(zip(tfs, wr_tf, trades_tf, pl_tf)):
    ax = axes[i]
    ax.set_facecolor('#1A243B')
    w = trades * wr / 100
    l = trades * (1 - wr / 100)
    bars = ax.bar(['Wins', 'Losses'], [w, l], color=['#00C853', '#FF4545'], width=0.6)
    ax.set_title(f'{tf} Timeframe\nWR: {wr}% | {trades} trades\nP/L: ${pl:,.0f}', color='white', fontsize=13, fontweight='bold', pad=10)
    ax.tick_params(colors='#8A93A6')
    for spine in ax.spines.values(): spine.set_color('#8A93A6')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, p: f'{x:,.0f}'))
    for bar in bars:
        ax.text(bar.get_x() + bar.get_width()/2., bar.get_height() + 30, f'{bar.get_height():,.0f}', ha='center', va='bottom', color='white', fontsize=11, fontweight='bold')
plt.tight_layout()
plt.savefig('m1m5_timeframes.png', dpi=150, facecolor='#0F1729', bbox_inches='tight')
plt.close()
slide.shapes.add_picture('m1m5_timeframes.png', Inches(1.5), Inches(1.2), Inches(10), Inches(5.8))

# ── Slide 6: Monthly P/L ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, Inches(0), Inches(0.2), Inches(13.333), Inches(0.7), "MONTHLY PROFIT & LOSS", 36, GOLD, True, PP_ALIGN.CENTER)
fig, ax = plt.subplots(figsize=(11, 4.5))
fig.patch.set_facecolor('#0F1729')
ax.set_facecolor('#0F1729')
colors = ['#00C853'] * len(monthly_pnl)
bars = ax.bar(months_labels, monthly_pnl, color=colors, width=0.5)
ax.set_xticks(range(len(months_labels)))
ax.set_xticklabels(months_labels, color='#8A93A6', fontsize=12)
ax.tick_params(axis='y', colors='#8A93A6')
ax.set_ylabel('Profit ($)', color='#8A93A6', fontsize=11)
ax.spines['bottom'].set_color('#8A93A6')
ax.spines['left'].set_color('#8A93A6')
ax.spines['top'].set_visible(False)
ax.spines['right'].set_visible(False)
ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, p: f'${x:,.0f}'))
for bar, val in zip(bars, monthly_pnl):
    ax.text(bar.get_x() + bar.get_width()/2., bar.get_height() + 500, f'${val:,.0f}', ha='center', va='bottom', color='white', fontsize=12, fontweight='bold')
plt.tight_layout()
plt.savefig('m1m5_monthly.png', dpi=150, facecolor='#0F1729', bbox_inches='tight')
plt.close()
slide.shapes.add_picture('m1m5_monthly.png', Inches(1.2), Inches(1.2), Inches(11), Inches(5.8))

# ── Slide 7: Settings ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, Inches(0), Inches(0.2), Inches(13.333), Inches(0.7), "OPTIMIZED SETTINGS", 36, GOLD, True, PP_ALIGN.CENTER)
settings = [
    ("FIBONACCI ENTRY ZONE", "0.500 - 0.786 (widened from 0.618-0.786)"),
    ("EMA PERIOD (TREND)", "30 (optimized from 50 - faster trend detection)"),
    ("TRAILING STOP START", "0.75 * ATR (triggers after 75% of ATR move)"),
    ("TRAILING STOP STEP", "0.2 * ATR (tighter than 0.3 - locks more profit)"),
    ("RISK PER TRADE", "4.0% of balance"),
    ("STOP LOSS", "1.0 * ATR (tight for scalping)"),
    ("TIMEFRAMES", "M1 + M5 (ultra-fast scalping)"),
    ("SYMBOLS", "XAUUSD, GBPUSD, AUDUSD"),
    ("ENTRY CONFIRMATION", "1 candle confirmation required"),
    ("MAX TRADES", "10 concurrent, 1 per symbol"),
]
y = Inches(1.1)
for i, (label, value) in enumerate(settings):
    bg = DARK if i % 2 == 0 else RGBColor(0x15, 0x1E, 0x33)
    add_box(slide, Inches(1), y, Inches(11.3), Inches(0.55), bg)
    add_text(slide, Inches(1.3), y + Inches(0.08), Inches(3.5), Inches(0.4), label, 14, GRAY, True)
    add_text(slide, Inches(5), y + Inches(0.08), Inches(7), Inches(0.4), value, 14, WHITE)
    y += Inches(0.58)

# ── Slide 8: Why It Works ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, Inches(0), Inches(0.2), Inches(13.333), Inches(0.7), "WHY THESE SETTINGS WORK", 36, GOLD, True, PP_ALIGN.CENTER)
reasons = [
    "M1+M5 timeframes generate 10,500+ trades in 3 months for high opportunity",
    "EMA 30 detects trend changes faster - critical for ultra-fast scalping",
    "Tighter trail step (0.2) captures quick moves without giving back profit",
    "Entry zone 0.500-0.786 catches retracements early on lower timeframes",
    "66.1% win rate with 10,500+ trades = massive compounding effect",
    "M5 timeframe (68% WR) outperforms M1 (64.6%) - slightly slower = cleaner signals",
    "XAUUSD leads with $45.5K profit - gold's volatility suits fast scalping",
]
y = Inches(1.2)
for i, reason in enumerate(reasons):
    add_box(slide, Inches(0.8), y, Inches(11.7), Inches(0.7), DARK if i % 2 == 0 else RGBColor(0x15, 0x1E, 0x33))
    add_text(slide, Inches(1.2), y + Inches(0.12), Inches(11), Inches(0.5), f"  {reason}", 15, WHITE)
    y += Inches(0.78)

# ── Save ──
output = "M1M5_Scalper_Presentation.pptx"
prs.save(output)
print(f"Saved: {output}")
