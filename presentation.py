"""
Kingade Scalper Bot - Creative PowerPoint Presentation
Generates a visually stunning .pptx with charts, metrics, and breakdowns.
"""

import sys
import os
import MetaTrader5 as mt5
import numpy as np
from datetime import datetime
from collections import defaultdict
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

import config
import backtest as bt


# ─── Color Palette ────────────────────────────────────────────────
DARK_BG = RGBColor(0x0D, 0x11, 0x17)
NAVY = RGBColor(0x16, 0x21, 0x3E)
BLUE = RGBColor(0x0F, 0x34, 0x60)
LIGHT_BLUE = RGBColor(0x1A, 0x73, 0xE8)
GOLD = RGBColor(0xF0, 0xA5, 0x00)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xE7, 0x4C, 0x3C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
MEDIUM_GRAY = RGBColor(0x99, 0x99, 0x99)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)


# ─── Helpers ──────────────────────────────────────────────────────
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=14, color=WHITE,
             bold=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_kpi_card(slide, left, top, label, value, color=WHITE, accent=GOLD):
    """Add a styled KPI card."""
    card_w = Inches(2.2)
    card_h = Inches(1.2)

    # Card background
    card = add_shape(slide, left, top, card_w, card_h, NAVY)
    card.shadow.inherit = False

    # Accent line at top
    accent_line = add_shape(slide, left, top, card_w, Pt(4), accent)

    # Label
    add_text(slide, left + Inches(0.15), top + Inches(0.2), card_w - Inches(0.3), Inches(0.35),
             label, font_size=9, color=MEDIUM_GRAY, bold=False)

    # Value
    add_text(slide, left + Inches(0.15), top + Inches(0.55), card_w - Inches(0.3), Inches(0.5),
             value, font_size=20, color=color, bold=True)


def add_table(slide, left, top, width, height, headers, rows, col_widths=None):
    """Add a styled table."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(9)
            paragraph.font.color.rgb = WHITE
            paragraph.font.bold = True
            paragraph.font.name = "Calibri"
            paragraph.alignment = PP_ALIGN.CENTER

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            bg = RGBColor(0x1A, 0x1A, 0x2E) if i % 2 == 0 else DARK_BG
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(8)
                paragraph.font.color.rgb = WHITE
                paragraph.font.name = "Calibri"
                paragraph.alignment = PP_ALIGN.CENTER

    return table_shape


# ─── Data Collection ──────────────────────────────────────────────
def collect_data():
    print("Running backtest for presentation data...")
    result = bt.run_backtest()
    return result


def build_weekly_data(trades):
    weekly = defaultdict(lambda: {
        "trades": 0, "wins": 0, "losses": 0, "pnl": 0.0,
        "start_date": None, "end_date": None
    })
    for t in trades:
        d = t.entry_time.date()
        iso_year, iso_week, _ = d.isocalendar()
        key = f"{iso_year}-W{iso_week:02d}"
        w = weekly[key]
        w["trades"] += 1
        w["pnl"] += t.profit
        if t.result == "win":
            w["wins"] += 1
        else:
            w["losses"] += 1
        if w["start_date"] is None or d < w["start_date"]:
            w["start_date"] = d
        if w["end_date"] is None or d > w["end_date"]:
            w["end_date"] = d
    return dict(sorted(weekly.items()))


def build_daily_data(trades):
    daily = defaultdict(lambda: {"trades": 0, "wins": 0, "losses": 0, "pnl": 0.0})
    for t in trades:
        day = t.entry_time.date()
        daily[day]["trades"] += 1
        daily[day]["pnl"] += t.profit
        if t.result == "win":
            daily[day]["wins"] += 1
        else:
            daily[day]["losses"] += 1
    return dict(sorted(daily.items()))


def build_monthly_data(trades):
    monthly = defaultdict(lambda: {"trades": 0, "wins": 0, "losses": 0, "pnl": 0.0})
    for t in trades:
        key = t.entry_time.strftime("%Y-%m")
        monthly[key]["trades"] += 1
        monthly[key]["pnl"] += t.profit
        if t.result == "win":
            monthly[key]["wins"] += 1
        else:
            monthly[key]["losses"] += 1
    return dict(sorted(monthly.items()))


# ─── Slide Builders ───────────────────────────────────────────────
def slide_title(prs):
    """Slide 1: Title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, DARK_BG)

    # Top accent bar
    add_shape(slide, Inches(0), Inches(0), Inches(13.33), Pt(6), GOLD)

    # Logo area - large "K" in gold circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.67), Inches(1.2), Inches(2), Inches(2))
    circle.fill.solid()
    circle.fill.fore_color.rgb = GOLD
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = "K"
    p.font.size = Pt(60)
    p.font.color.rgb = DARK_BG
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(8)

    # Title
    add_text(slide, Inches(1), Inches(3.5), Inches(11.33), Inches(1),
             "KINGADE SCALPER BOT", font_size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Subtitle
    add_text(slide, Inches(1), Inches(4.5), Inches(11.33), Inches(0.6),
             "Automated Fibonacci Retracement Scalping System", font_size=18, color=GOLD, align=PP_ALIGN.CENTER)

    # Date
    add_text(slide, Inches(1), Inches(5.3), Inches(11.33), Inches(0.4),
             f"Performance Report | {datetime.now().strftime('%B %Y')}", font_size=12, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)

    # Bottom accent bar
    add_shape(slide, Inches(0), Inches(7.45), Inches(13.33), Pt(4), GOLD)

    # Footer text
    add_text(slide, Inches(1), Inches(6.8), Inches(11.33), Inches(0.3),
             "XAUUSD | M1  M5  M15  M30 | 8% Risk | 548 Trades", font_size=10, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)


def slide_executive_summary(prs, result):
    """Slide 2: Executive Summary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # Title bar
    add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.9), NAVY)
    add_text(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
             "EXECUTIVE SUMMARY", font_size=24, color=GOLD, bold=True)

    # Row 1 KPIs
    kpi_y = Inches(1.2)
    add_kpi_card(slide, Inches(0.5), kpi_y, "INITIAL BALANCE", f"${result['initial_balance']:,.0f}", WHITE, BLUE)
    add_kpi_card(slide, Inches(2.9), kpi_y, "FINAL BALANCE", f"${result['final_balance']:,.0f}", GREEN, GREEN)
    add_kpi_card(slide, Inches(5.3), kpi_y, "NET PROFIT", f"${result['total_pnl']:,.0f} ({result['total_pnl_pct']:+.0f}%)", GOLD, GOLD)
    add_kpi_card(slide, Inches(7.7), kpi_y, "WIN RATE", f"{result['win_rate']:.1f}%", GREEN, GREEN)
    add_kpi_card(slide, Inches(10.1), kpi_y, "PROFIT FACTOR", f"{result['profit_factor']:.2f}", GOLD, GOLD)

    # Row 2 KPIs
    kpi_y2 = Inches(2.6)
    add_kpi_card(slide, Inches(0.5), kpi_y2, "TOTAL TRADES", str(result['total_trades']), WHITE, BLUE)
    add_kpi_card(slide, Inches(2.9), kpi_y2, "WINNING TRADES", str(result['wins']), GREEN, GREEN)
    add_kpi_card(slide, Inches(5.3), kpi_y2, "LOSING TRADES", str(result['losses']), RED, RED)
    add_kpi_card(slide, Inches(7.7), kpi_y2, "AVG WIN", f"${result['avg_win']:,.2f}", GREEN, GREEN)
    add_kpi_card(slide, Inches(10.1), kpi_y2, "AVG LOSS", f"${result['avg_loss']:,.2f}", RED, RED)

    # Row 3 KPIs
    kpi_y3 = Inches(4.0)
    add_kpi_card(slide, Inches(0.5), kpi_y3, "SHARPE RATIO", f"{result['sharpe']:.2f}", GOLD, GOLD)
    add_kpi_card(slide, Inches(2.9), kpi_y3, "MAX DRAWDOWN", f"${result['max_dd']:,.0f} ({result['max_dd_pct']:.1f}%)", RED, RED)
    add_kpi_card(slide, Inches(5.3), kpi_y3, "AVG R:R", f"{result['avg_rr']:.2f}", WHITE, BLUE)
    add_kpi_card(slide, Inches(7.7), kpi_y3, "EXPECTANCY", f"${result['expectancy']:,.2f}/trade", GOLD, GOLD)
    add_kpi_card(slide, Inches(10.1), kpi_y3, "AVG BARS HELD", f"{result['avg_bars_held']:.1f}", WHITE, BLUE)

    # Bottom insight box
    add_shape(slide, Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.5), NAVY, GOLD)
    insight_lines = [
        f"From $1,000 to ${result['final_balance']:,.0f} = {result['total_pnl_pct']:+.0f}% return",
        f"Avg ${result['total_pnl']/13:,.0f}/month over 13 months | {result['wins']}W / {result['losses']}L = {result['win_rate']:.1f}% win rate",
        f"Largest Win: ${result['largest_win']:,.2f} | Largest Loss: ${result['largest_loss']:,.2f} | Max DD: {result['max_dd_pct']:.1f}%"
    ]
    add_text(slide, Inches(0.8), Inches(5.75), Inches(11.7), Inches(1.2),
             "\n".join(insight_lines), font_size=11, color=WHITE, align=PP_ALIGN.CENTER)


def slide_monthly_performance(prs, result):
    """Slide 3: Monthly Performance with chart"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.9), NAVY)
    add_text(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
             "MONTHLY PERFORMANCE", font_size=24, color=GOLD, bold=True)

    monthly = build_monthly_data(result["trades"])

    # Monthly bar chart using built-in chart
    chart_data = __import__('pptx.chart.data', fromlist=['CategoryChartData']).CategoryChartData()
    months = list(monthly.keys())
    pnls = [monthly[m]["pnl"] for m in months]
    chart_data.categories = months
    chart_data.add_series("P/L ($)", pnls)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.5), Inches(1.1), Inches(8), Inches(3.5),
        chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.style = 2

    # Style the chart
    plot = chart.plots[0]
    plot.gap_width = 80
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = GOLD

    # Chart area
    chart.chart_style = 2
    chart_format = chart.element
    chart_area = chart_format

    # Value axis
    value_axis = chart.value_axis
    value_axis.has_title = True
    value_axis.axis_title.text_frame.paragraphs[0].text = "P/L ($)"
    value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(9)
    value_axis.axis_title.text_frame.paragraphs[0].font.color.rgb = MEDIUM_GRAY
    value_axis.major_gridlines.format.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    value_axis.format.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    value_axis.tick_labels.font.size = Pt(8)
    value_axis.tick_labels.font.color.rgb = MEDIUM_GRAY

    # Category axis
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(7)
    cat_axis.tick_labels.font.color.rgb = MEDIUM_GRAY
    cat_axis.format.line.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # Monthly table on the right
    headers = ["MONTH", "TRADES", "W", "L", "WR%", "P/L"]
    rows = []
    cum = 0
    for m, data in monthly.items():
        cum += data["pnl"]
        wr = (data["wins"] / data["trades"] * 100) if data["trades"] else 0
        rows.append([m, str(data['trades']), str(data['wins']), str(data['losses']),
                     f"{wr:.0f}%", f"${data['pnl']:+,.0f}"])

    add_table(slide, Inches(8.8), Inches(1.1), Inches(4.2), Inches(3.5),
              headers, rows)

    # Cumulative P/L row at bottom
    cum_data = []
    cum = 0
    for m in months:
        cum += monthly[m]["pnl"]
        cum_data.append(cum)

    # Bottom insight
    avg_monthly = result['total_pnl'] / len(months) if months else 0
    best_month = max(monthly.items(), key=lambda x: x[1]["pnl"])
    worst_month = min(monthly.items(), key=lambda x: x[1]["pnl"])

    add_shape(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.2), NAVY, GOLD)
    add_text(slide, Inches(0.8), Inches(5.1), Inches(11.7), Inches(2.0),
             f"Avg Monthly: ${avg_monthly:,.0f}     |     "
             f"Best Month: {best_month[0]} (${best_month[1]['pnl']:+,.0f})     |     "
             f"Worst Month: {worst_month[0]} (${worst_month[1]['pnl']:+,.0f})\n\n"
             f"Cumulative P/L: ${result['total_pnl']:,.0f}     |     "
             f"Positive Months: {sum(1 for m in monthly.values() if m['pnl'] > 0)}/{len(months)}     |     "
             f"Profit Factor: {result['profit_factor']:.2f}",
             font_size=11, color=WHITE, align=PP_ALIGN.CENTER)


def slide_weekly_breakdown(prs, result, page=0, per_page=6):
    """Slide 4+: Weekly breakdown pages"""
    weekly = build_weekly_data(result["trades"])
    weeks = list(weekly.items())
    start = page * per_page
    end = min(start + per_page, len(weeks))

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.9), NAVY)
    title = f"WEEKLY BREAKDOWN" + (f" ({page+1})" if len(weeks) > per_page else "")
    add_text(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
             title, font_size=24, color=GOLD, bold=True)

    # Two columns of 3 weeks each
    col_w = Inches(6.0)
    card_h = Inches(1.8)

    cum_before = sum(weekly[w]["pnl"] for w in list(weekly.keys())[:start])

    for idx in range(start, end):
        week_key, data = weeks[idx]
        col = (idx - start) % 2
        row = (idx - start) // 2

        x = Inches(0.4) + col * Inches(6.4)
        y = Inches(1.1) + row * Inches(2.0)

        wr = (data["wins"] / data["trades"] * 100) if data["trades"] else 0
        cum_before += data["pnl"]
        pnl_color = GREEN if data["pnl"] >= 0 else RED

        # Card background
        card = add_shape(slide, x, y, col_w, card_h, NAVY)

        # Week header
        add_text(slide, x + Inches(0.15), y + Inches(0.05), col_w - Inches(0.3), Inches(0.3),
                 f"{week_key}  |  {data['start_date'].strftime('%d %b')} - {data['end_date'].strftime('%d %b %Y')}",
                 font_size=9, color=GOLD, bold=True)

        # Stats row 1
        stats1 = f"Trades: {data['trades']}   |   Wins: {data['wins']}   |   Losses: {data['losses']}   |   WR: {wr:.0f}%"
        add_text(slide, x + Inches(0.15), y + Inches(0.4), col_w - Inches(0.3), Inches(0.3),
                 stats1, font_size=9, color=WHITE)

        # Stats row 2 - P/L big
        add_text(slide, x + Inches(0.15), y + Inches(0.75), Inches(2.5), Inches(0.5),
                 f"${data['pnl']:+,.2f}", font_size=24, color=pnl_color, bold=True)

        # Cumulative
        add_text(slide, x + Inches(3.5), y + Inches(0.85), Inches(2.3), Inches(0.3),
                 f"Cumulative: ${cum_before:,.0f}", font_size=10, color=MEDIUM_GRAY, align=PP_ALIGN.RIGHT)

        # Win/Loss bar
        bar_w = col_w - Inches(0.3)
        bar_h = Inches(0.2)
        bar_y = y + Inches(1.35)
        if data['trades'] > 0:
            win_w = int(bar_w * data['wins'] / data['trades'])
            loss_w = bar_w - win_w
            if win_w > 0:
                win_bar = add_shape(slide, x + Inches(0.15), bar_y, win_w, bar_h, GREEN)
            if loss_w > 0:
                loss_bar = add_shape(slide, x + Inches(0.15) + win_w, bar_y, loss_w, bar_h, RED)

        # Label
        add_text(slide, x + Inches(0.15), bar_y + bar_h + Inches(0.02), bar_w, Inches(0.2),
                 f"{'W' * data['wins']}{'L' * data['losses']}", font_size=6, color=MEDIUM_GRAY)


def slide_by_timeframe(prs, result):
    """Slide 5: Performance by Timeframe"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.9), NAVY)
    add_text(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
             "PERFORMANCE BY TIMEFRAME", font_size=24, color=GOLD, bold=True)

    tf_stats = result["tf_stats"]
    tf_order = ["M1", "M5", "M15", "M30"]

    # Timeframe cards - large visual cards
    card_w = Inches(2.8)
    card_h = Inches(4.5)
    start_x = Inches(0.5)
    gap = Inches(0.35)

    for i, tf in enumerate(tf_order):
        if tf not in tf_stats:
            continue
        s = tf_stats[tf]
        wr = (s["wins"] / s["trades"] * 100) if s["trades"] else 0

        x = start_x + i * (card_w + gap)
        y = Inches(1.2)

        # Card
        add_shape(slide, x, y, card_w, card_h, NAVY)

        # TF name - large
        add_text(slide, x, y + Inches(0.2), card_w, Inches(0.8),
                 tf, font_size=36, color=GOLD, bold=True, align=PP_ALIGN.CENTER)

        # Trades
        add_text(slide, x, y + Inches(1.1), card_w, Inches(0.3),
                 f"{s['trades']} Trades", font_size=14, color=WHITE, align=PP_ALIGN.CENTER)

        # Win Rate - big number
        wr_color = GREEN if wr >= 70 else GOLD if wr >= 60 else RED
        add_text(slide, x, y + Inches(1.6), card_w, Inches(0.7),
                 f"{wr:.1f}%", font_size=32, color=wr_color, bold=True, align=PP_ALIGN.CENTER)

        add_text(slide, x, y + Inches(2.2), card_w, Inches(0.3),
                 "Win Rate", font_size=10, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)

        # Wins / Losses
        add_text(slide, x, y + Inches(2.7), card_w, Inches(0.4),
                 f"W: {s['wins']}  |  L: {s['losses']}", font_size=12, color=WHITE, align=PP_ALIGN.CENTER)

        # P/L - big
        pnl_color = GREEN if s["pnl"] > 0 else RED
        add_text(slide, x, y + Inches(3.2), card_w, Inches(0.6),
                 f"${s['pnl']:,.0f}", font_size=22, color=pnl_color, bold=True, align=PP_ALIGN.CENTER)

        # P/L as % of total
        pct = (s["pnl"] / result["total_pnl"] * 100) if result["total_pnl"] else 0
        add_text(slide, x, y + Inches(3.8), card_w, Inches(0.3),
                 f"{pct:.0f}% of total P/L", font_size=9, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)

    # Bottom insight
    best_tf = max(tf_stats.items(), key=lambda x: x[1]["pnl"])
    add_shape(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.2), NAVY, GOLD)
    add_text(slide, Inches(0.8), Inches(6.15), Inches(11.7), Inches(0.9),
             f"Best performing timeframe: {best_tf[0]} with ${best_tf[1]['pnl']:,.0f} profit "
             f"({(best_tf[1]['pnl']/result['total_pnl']*100):.0f}% of total) | "
             f"M30 is the most consistent with highest win rate",
             font_size=11, color=WHITE, align=PP_ALIGN.CENTER)


def slide_drawdown_analysis(prs, result):
    """Slide 6: Drawdown Analysis"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.9), NAVY)
    add_text(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
             "DRAWDOWN & RISK ANALYSIS", font_size=24, color=GOLD, bold=True)

    # Left: Drawdown metrics
    add_shape(slide, Inches(0.5), Inches(1.2), Inches(5.5), Inches(3.0), NAVY)

    dd_items = [
        ("Max Drawdown", f"${result['max_dd']:,.2f}", RED),
        ("Max Drawdown %", f"{result['max_dd_pct']:.2f}%", RED),
        ("Sharpe Ratio", f"{result['sharpe']:.2f}", GOLD),
        ("Recovery Factor", f"{result['recovery']:.2f}", GREEN),
        ("Calmar Ratio", f"{result['calmar']:.2f}", GREEN),
    ]

    for i, (label, value, color) in enumerate(dd_items):
        y = Inches(1.4) + i * Inches(0.5)
        add_text(slide, Inches(0.8), y, Inches(2.5), Inches(0.35),
                 label, font_size=11, color=MEDIUM_GRAY)
        add_text(slide, Inches(3.5), y, Inches(2.2), Inches(0.35),
                 value, font_size=14, color=color, bold=True, align=PP_ALIGN.RIGHT)

    # Right: Risk visualization
    add_shape(slide, Inches(6.3), Inches(1.2), Inches(6.5), Inches(3.0), NAVY)

    # Risk meter visual
    meter_x = Inches(6.8)
    meter_y = Inches(1.5)
    meter_w = Inches(5.5)
    meter_h = Inches(0.6)

    # Background bar
    add_shape(slide, meter_x, meter_y, meter_w, meter_h, RGBColor(0x33, 0x33, 0x33))

    # Fill based on drawdown (20% = red zone)
    dd_fill_w = int(meter_w * min(result["max_dd_pct"] / 100, 1.0))
    dd_color = RED if result["max_dd_pct"] > 20 else GOLD if result["max_dd_pct"] > 10 else GREEN
    add_shape(slide, meter_x, meter_y, dd_fill_w, meter_h, dd_color)

    add_text(slide, meter_x, meter_y + meter_h + Inches(0.05), meter_w, Inches(0.3),
             f"Drawdown: {result['max_dd_pct']:.1f}%   |   Risk Level: {'HIGH' if result['max_dd_pct'] > 20 else 'MODERATE' if result['max_dd_pct'] > 10 else 'LOW'}",
             font_size=9, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)

    # Risk/Reward visualization
    rr_items = [
        ("Risk Per Trade", f"{config.RISK_PERCENT}%", BLUE),
        ("Avg Win", f"${result['avg_win']:,.2f}", GREEN),
        ("Avg Loss", f"${result['avg_loss']:,.2f}", RED),
        ("Avg R:R", f"{result['avg_rr']:.2f}", GOLD),
        ("Expectancy", f"${result['expectancy']:,.2f}", GREEN),
    ]

    for i, (label, value, color) in enumerate(rr_items):
        y = Inches(2.6) + i * Inches(0.35)
        add_text(slide, Inches(6.6), y, Inches(2.5), Inches(0.3),
                 label, font_size=10, color=MEDIUM_GRAY)
        add_text(slide, Inches(9.5), y, Inches(3.0), Inches(0.3),
                 value, font_size=12, color=color, bold=True, align=PP_ALIGN.RIGHT)

    # Bottom: Equity curve simulation (simplified)
    add_shape(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.7), NAVY, GOLD)

    # Build equity curve data
    bal = result["initial_balance"]
    eq_points = [(0, bal)]
    for t in sorted(result["trades"], key=lambda x: x.entry_time):
        bal += t.profit
        eq_points.append((len(eq_points), round(bal, 2)))

    # Draw simplified equity curve using shapes
    chart_w = Inches(11.5)
    chart_h = Inches(2.0)
    chart_x = Inches(0.9)
    chart_y = Inches(4.8)

    vals = [p[1] for p in eq_points]
    mn, mx = min(vals), max(vals)
    range_val = mx - mn if mx != mn else 1

    # Draw as a series of thin vertical bars
    n_bars = min(len(vals), 200)
    step = max(1, len(vals) // n_bars)
    sampled = vals[::step]

    bar_w = max(int(chart_w / n_bars), Emu(1000))
    for i, v in enumerate(sampled):
        x = chart_x + int(chart_w * i / n_bars)
        h = int(chart_h * (v - mn) / range_val) if range_val > 0 else int(chart_h / 2)
        bar_color = GREEN if v >= result["initial_balance"] else RED
        add_shape(slide, x, chart_y + chart_h - h, bar_w, h, bar_color)

    # Labels
    add_text(slide, chart_x, chart_y - Inches(0.25), Inches(3), Inches(0.2),
             f"Starting: ${result['initial_balance']:,.0f}", font_size=8, color=MEDIUM_GRAY)
    add_text(slide, chart_x + chart_w - Inches(3), chart_y - Inches(0.25), Inches(3), Inches(0.2),
             f"Peak: ${mx:,.0f}  |  End: ${result['final_balance']:,.0f}", font_size=8, color=MEDIUM_GRAY, align=PP_ALIGN.RIGHT)


def slide_trade_log_summary(prs, result):
    """Slide 7: Trade Log Summary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.9), NAVY)
    add_text(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
             "TRADE LOG SUMMARY", font_size=24, color=GOLD, bold=True)

    trades = sorted(result["trades"], key=lambda x: x.entry_time)

    # Stats cards
    # Streak analysis
    max_win_streak = 0
    max_loss_streak = 0
    current_streak = 0
    streak_type = None
    for t in trades:
        if t.result == "win":
            if streak_type == "win":
                current_streak += 1
            else:
                current_streak = 1
                streak_type = "win"
            max_win_streak = max(max_win_streak, current_streak)
        else:
            if streak_type == "loss":
                current_streak += 1
            else:
                current_streak = 1
                streak_type = "loss"
            max_loss_streak = max(max_loss_streak, current_streak)

    # Direction split
    buy_trades = [t for t in trades if t.direction == "buy"]
    sell_trades = [t for t in trades if t.direction == "sell"]
    buy_wins = sum(1 for t in buy_trades if t.result == "win")
    sell_wins = sum(1 for t in sell_trades if t.result == "win")

    # KPIs
    kpi_y = Inches(1.2)
    add_kpi_card(slide, Inches(0.5), kpi_y, "MAX WIN STREAK", str(max_win_streak), GREEN, GREEN)
    add_kpi_card(slide, Inches(2.9), kpi_y, "MAX LOSS STREAK", str(max_loss_streak), RED, RED)
    add_kpi_card(slide, Inches(5.3), kpi_y, "BUY TRADES", f"{len(buy_trades)} ({buy_wins/len(buy_trades)*100:.0f}% WR)" if buy_trades else "0", WHITE, BLUE)
    add_kpi_card(slide, Inches(7.7), kpi_y, "SELL TRADES", f"{len(sell_trades)} ({sell_wins/len(sell_trades)*100:.0f}% WR)" if sell_trades else "0", WHITE, BLUE)
    add_kpi_card(slide, Inches(10.1), kpi_y, "TOTAL P/L", f"${result['total_pnl']:,.0f}", GOLD, GOLD)

    # Direction breakdown table
    add_shape(slide, Inches(0.5), Inches(2.7), Inches(6.0), Inches(4.3), NAVY)
    add_text(slide, Inches(0.7), Inches(2.85), Inches(5.5), Inches(0.3),
             "DIRECTION BREAKDOWN", font_size=12, color=GOLD, bold=True)

    dir_headers = ["", "TRADES", "WINS", "LOSSES", "WR%", "P/L"]
    dir_rows = [
        ["BUY", str(len(buy_trades)), str(buy_wins), str(len(buy_trades)-buy_wins),
         f"{buy_wins/len(buy_trades)*100:.1f}%" if buy_trades else "0%",
         f"${sum(t.profit for t in buy_trades):+,.0f}"],
        ["SELL", str(len(sell_trades)), str(sell_wins), str(len(sell_trades)-sell_wins),
         f"{sell_wins/len(sell_trades)*100:.1f}%" if sell_trades else "0%",
         f"${sum(t.profit for t in sell_trades):+,.0f}"],
    ]
    add_table(slide, Inches(0.7), Inches(3.3), Inches(5.5), Inches(1.2), dir_headers, dir_rows)

    # Recent trades (last 15)
    add_shape(slide, Inches(6.8), Inches(2.7), Inches(6.0), Inches(4.3), NAVY)
    add_text(slide, Inches(7.0), Inches(2.85), Inches(5.5), Inches(0.3),
             "LAST 15 TRADES", font_size=12, color=GOLD, bold=True)

    recent_headers = ["DATE", "TF", "DIR", "P/L", "RESULT"]
    recent_rows = []
    for t in trades[-15:]:
        recent_rows.append([
            t.entry_time.strftime("%d/%m %H:%M"),
            bt._tf_name(t.timeframe),
            t.direction[0].upper(),
            f"${t.profit:+,.2f}",
            "WIN" if t.result == "win" else "LOSS"
        ])

    add_table(slide, Inches(7.0), Inches(3.3), Inches(5.5), Inches(3.5),
              recent_headers, recent_rows)


def slide_closing(prs, result):
    """Slide 8: Closing"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # Top accent
    add_shape(slide, Inches(0), Inches(0), Inches(13.33), Pt(6), GOLD)

    # Logo
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.67), Inches(1.5), Inches(2), Inches(2))
    circle.fill.solid()
    circle.fill.fore_color.rgb = GOLD
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = "K"
    p.font.size = Pt(60)
    p.font.color.rgb = DARK_BG
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(8)

    add_text(slide, Inches(1), Inches(3.8), Inches(11.33), Inches(0.8),
             "KINGADE SCALPER BOT", font_size=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(1), Inches(4.6), Inches(11.33), Inches(0.5),
             f"${result['initial_balance']:,.0f}  →  ${result['final_balance']:,.0f}  |  "
             f"{result['total_pnl_pct']:+.0f}% Return  |  "
             f"{result['win_rate']:.0f}% Win Rate  |  "
             f"Sharpe {result['sharpe']:.1f}",
             font_size=16, color=GOLD, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(1), Inches(5.5), Inches(11.33), Inches(0.5),
             "Automated Fibonacci Scalping | XAUUSD | M1-M30",
             font_size=14, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)

    # Bottom accent
    add_shape(slide, Inches(0), Inches(7.45), Inches(13.33), Pt(4), GOLD)


# ─── Main ─────────────────────────────────────────────────────────
def create_presentation():
    result = collect_data()

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    print("Building slides...")
    slide_title(prs)
    slide_executive_summary(prs, result)
    slide_monthly_performance(prs, result)

    # Weekly breakdown - multiple pages
    weekly = build_weekly_data(result["trades"])
    n_weeks = len(weekly)
    per_page = 6
    n_pages = (n_weeks + per_page - 1) // per_page
    for page in range(n_pages):
        slide_weekly_breakdown(prs, result, page=page, per_page=per_page)

    slide_by_timeframe(prs, result)
    slide_drawdown_analysis(prs, result)
    slide_trade_log_summary(prs, result)
    slide_closing(prs, result)

    # Save
    report_dir = os.path.dirname(os.path.abspath(__file__))
    timestamp = datetime.now().strftime("%Y%m%d_%H%M")
    output_path = os.path.join(report_dir, f"Kingade_Scalper_Bot_{timestamp}.pptx")
    prs.save(output_path)
    print(f"\nPresentation saved: {output_path}")
    return output_path


if __name__ == "__main__":
    create_presentation()
    mt5.shutdown()
