# Copyright (c) 2026 Kingade Forex. All rights reserved.
# This software is licensed intellectual property.
# Unauthorized copying, modification, distribution, or use is strictly prohibited.
# A valid license key (KNG-XXXX-XXXX-XXXX) is required to run this bot.
# Purchase at: https://sellix.io/kingadebot
"""Daily report generator for Kingade Scalper Bot.
Generates PDF and PPTX daily reports and sends them via Telegram.
"""

import os
import logging
import MetaTrader5 as mt5
import numpy as np
from datetime import datetime, timedelta, timezone
from collections import defaultdict
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import mm
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import (
    SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, HRFlowable
)
from reportlab.graphics.shapes import Drawing, Rect
from reportlab.graphics.charts.lineplots import LinePlot
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE

import config
import telegram_notifier as tg

# Colors
DARK_BG = RGBColor(0x0D, 0x11, 0x17)
NAVY = RGBColor(0x16, 0x21, 0x3E)
BLUE = RGBColor(0x0F, 0x34, 0x60)
GOLD = RGBColor(0xF0, 0xA5, 0x00)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xE7, 0x4C, 0x3C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MEDIUM_GRAY = RGBColor(0x99, 0x99, 0x99)

logger = logging.getLogger(__name__)


def get_today_trades():
    """Get all trades opened today from MT5 deal history.
    MT5 timestamps are UTC, so the day window must be UTC-based."""
    now = datetime.now(timezone.utc)
    today_start = datetime(now.year, now.month, now.day, tzinfo=timezone.utc)

    from_ts = int(today_start.timestamp())
    to_ts = int(now.timestamp())

    deals = mt5.history_deals_get(from_ts, to_ts)
    if deals is None:
        return []

    bot_deals = [d for d in deals if d.magic == config.MAGIC_NUMBER]

    trades = []
    for d in bot_deals:
        trades.append({
            "ticket": d.order,
            "time": datetime.fromtimestamp(d.time),
            "symbol": d.symbol,
            "type": "BUY" if d.type == mt5.DEAL_TYPE_BUY else "SELL",
            "volume": d.volume,
            "price": d.price,
            "profit": d.profit,
            "swap": d.swap,
            "commission": d.commission,
            "net_profit": d.profit + d.swap + d.commission,
            "comment": d.comment,
        })

    return trades


def get_account_snapshot():
    """Get current account info."""
    info = mt5.account_info()
    if info is None:
        return {}
    return {
        "balance": info.balance,
        "equity": info.equity,
        "margin": info.margin,
        "free_margin": info.margin_free,
        "profit": info.profit,
    }


# ─── PDF Report ───────────────────────────────────────────────────
def generate_pdf(trades, account, output_path):
    """Generate daily PDF report."""
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle("Title2", parent=styles["Title"], fontSize=18, textColor=colors.HexColor("#1a1a2e")))
    styles.add(ParagraphStyle("Sub", parent=styles["Normal"], fontSize=10, textColor=colors.grey))

    doc = SimpleDocTemplate(output_path, pagesize=A4, leftMargin=15*mm, rightMargin=15*mm, topMargin=15*mm, bottomMargin=15*mm)
    elements = []

    now = datetime.now()
    date_str = now.strftime("%d %B %Y")

    # Title
    elements.append(Paragraph("Kingade Scalper Bot", styles["Title2"]))
    elements.append(Paragraph(f"Daily Report - {date_str}", styles["Sub"]))
    elements.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#16213e")))
    elements.append(Spacer(1, 8))

    # Account Summary
    wins = [t for t in trades if t["net_profit"] >= 0]
    losses = [t for t in trades if t["net_profit"] < 0]
    total_pnl = sum(t["net_profit"] for t in trades)

    summary = [
        ["DAILY SUMMARY", "", "", ""],
        ["Date", date_str, "Total Trades", str(len(trades))],
        ["Balance", f"${account.get('balance', 0):,.2f}", "Equity", f"${account.get('equity', 0):,.2f}"],
        ["Daily P/L", f"${total_pnl:+,.2f}", "Win Rate", f"{len(wins)/len(trades)*100:.0f}%" if trades else "0%"],
        ["Wins", str(len(wins)), "Losses", str(len(losses))],
        ["Avg Win", f"${np.mean([t['net_profit'] for t in wins]):+,.2f}" if wins else "$0", "Avg Loss", f"${np.mean([t['net_profit'] for t in losses]):+,.2f}" if losses else "$0"],
    ]

    t = Table(summary, colWidths=[100, 140, 100, 140])
    t.setStyle(TableStyle([
        ("SPAN", (0, 0), (3, 0)),
        ("BACKGROUND", (0, 0), (3, 0), colors.HexColor("#16213e")),
        ("TEXTCOLOR", (0, 0), (3, 0), colors.white),
        ("FONTNAME", (0, 0), (3, 0), "Helvetica-Bold"),
        ("FONTSIZE", (0, 0), (3, 0), 11),
        ("ALIGN", (0, 0), (3, 0), "CENTER"),
        ("FONTNAME", (0, 1), (0, -1), "Helvetica-Bold"),
        ("FONTNAME", (2, 1), (2, -1), "Helvetica-Bold"),
        ("FONTSIZE", (0, 1), (-1, -1), 9),
        ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#cccccc")),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f5f5f5")]),
        ("TOPPADDING", (0, 0), (-1, -1), 4),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
    ]))
    elements.append(t)
    elements.append(Spacer(1, 12))

    # Trade Log
    if trades:
        elements.append(Paragraph("Trade Log", styles["Heading2"]))
        elements.append(HRFlowable(width="100%", thickness=0.5, color=colors.HexColor("#cccccc")))
        elements.append(Spacer(1, 4))

        log_header = ["#", "TIME", "SYMBOL", "TYPE", "ENTRY", "VOLUME", "P/L", "RESULT"]
        log_rows = [log_header]
        for i, tr in enumerate(trades, 1):
            result = "WIN" if tr["net_profit"] >= 0 else "LOSS"
            log_rows.append([
                str(i),
                tr["time"].strftime("%H:%M"),
                tr["symbol"],
                tr["type"],
                f"{tr['price']:.2f}",
                f"{tr['volume']:.2f}",
                f"${tr['net_profit']:+,.2f}",
                result,
            ])

        t2 = Table(log_rows, colWidths=[30, 55, 60, 40, 70, 55, 70, 50])
        t2.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#0f3460")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, -1), 8),
            ("ALIGN", (0, 0), (-1, -1), "CENTER"),
            ("GRID", (0, 0), (-1, -1), 0.3, colors.HexColor("#dddddd")),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f8f8f8")]),
            ("TOPPADDING", (0, 0), (-1, -1), 3),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
        ]))
        elements.append(t2)
    else:
        elements.append(Paragraph("No trades today", styles["Heading3"]))

    doc.build(elements)
    return output_path


# ─── PPTX Report ──────────────────────────────────────────────────
def _add_title_bar(slide, text):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(22)
    p.font.color.rgb = GOLD
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5)
    tf.vertical_anchor = 1  # middle


def _add_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG


def _add_kpi(slide, x, y, label, value, color=WHITE, accent=GOLD):
    w, h = Inches(2.2), Inches(1.2)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    # accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = accent
    line.line.fill.background()
    # label
    tb1 = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.2), w - Inches(0.3), Inches(0.3))
    p1 = tb1.text_frame.paragraphs[0]
    p1.text = label
    p1.font.size = Pt(9)
    p1.font.color.rgb = MEDIUM_GRAY
    p1.font.name = "Calibri"
    # value
    tb2 = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.55), w - Inches(0.3), Inches(0.5))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = value
    p2.font.size = Pt(18)
    p2.font.color.rgb = color
    p2.font.bold = True
    p2.font.name = "Calibri"


def generate_pptx(trades, account, output_path):
    """Generate daily PPTX report."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    now = datetime.now()
    date_str = now.strftime("%d %B %Y")

    wins = [t for t in trades if t["net_profit"] >= 0]
    losses = [t for t in trades if t["net_profit"] < 0]
    total_pnl = sum(t["net_profit"] for t in trades)

    # ─── Slide 1: Overview ────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _add_title_bar(slide, f"KINGADE SCALPER BOT - Daily Report - {date_str}")

    y = Inches(1.2)
    _add_kpi(slide, Inches(0.5), y, "BALANCE", f"${account.get('balance', 0):,.2f}", WHITE, BLUE)
    _add_kpi(slide, Inches(2.9), y, "EQUITY", f"${account.get('equity', 0):,.2f}", WHITE, BLUE)
    _add_kpi(slide, Inches(5.3), y, "DAILY P/L", f"${total_pnl:+,.2f}", GREEN if total_pnl >= 0 else RED, GREEN if total_pnl >= 0 else RED)
    _add_kpi(slide, Inches(7.7), y, "TOTAL TRADES", str(len(trades)), WHITE, BLUE)
    _add_kpi(slide, Inches(10.1), y, "WIN RATE", f"{len(wins)/len(trades)*100:.0f}%" if trades else "0%", GREEN if trades and len(wins)/len(trades) >= 0.5 else RED, GREEN if trades and len(wins)/len(trades) >= 0.5 else RED)

    y2 = Inches(2.6)
    _add_kpi(slide, Inches(0.5), y2, "WINS", str(len(wins)), GREEN, GREEN)
    _add_kpi(slide, Inches(2.9), y2, "LOSSES", str(len(losses)), RED, RED)
    avg_win = np.mean([t["net_profit"] for t in wins]) if wins else 0
    avg_loss = np.mean([t["net_profit"] for t in losses]) if losses else 0
    _add_kpi(slide, Inches(5.3), y2, "AVG WIN", f"${avg_win:+,.2f}", GREEN, GREEN)
    _add_kpi(slide, Inches(7.7), y2, "AVG LOSS", f"${avg_loss:+,.2f}", RED, RED)
    gross_profit = sum(t["net_profit"] for t in wins) if wins else 0
    gross_loss = abs(sum(t["net_profit"] for t in losses)) if losses else 0
    pf = gross_profit / gross_loss if gross_loss > 0 else 0
    _add_kpi(slide, Inches(10.1), y2, "PROFIT FACTOR", f"{pf:.2f}", GOLD, GOLD)

    # ─── Slide 2: Trade Log ───────────────────────────────────────
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide2)
    _add_title_bar(slide2, f"TRADE LOG - {date_str}")

    if trades:
        n_rows = len(trades) + 1
        tbl = slide2.shapes.add_table(n_rows, 8, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.5))
        table = tbl.table

        headers = ["#", "TIME", "SYMBOL", "TYPE", "ENTRY", "VOLUME", "P/L", "RESULT"]
        for j, h in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = WHITE
                p.font.bold = True
                p.font.name = "Calibri"
                p.alignment = PP_ALIGN.CENTER

        for i, tr in enumerate(trades, 1):
            row_data = [
                str(i), tr["time"].strftime("%H:%M"), tr["symbol"], tr["type"],
                f"{tr['price']:.2f}", f"{tr['volume']:.2f}",
                f"${tr['net_profit']:+,.2f}", "WIN" if tr["net_profit"] >= 0 else "LOSS",
            ]
            for j, val in enumerate(row_data):
                cell = table.cell(i, j)
                cell.text = val
                bg = RGBColor(0x1A, 0x1A, 0x2E) if i % 2 == 0 else DARK_BG
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(9)
                    p.font.color.rgb = WHITE
                    p.font.name = "Calibri"
                    p.alignment = PP_ALIGN.CENTER
                    if j == 6:  # P/L column
                        p.font.color.rgb = GREEN if tr["net_profit"] >= 0 else RED
                        p.font.bold = True
                    if j == 7:  # Result column
                        p.font.color.rgb = GREEN if tr["net_profit"] >= 0 else RED

        # Column widths
        widths = [Inches(0.5), Inches(1.2), Inches(1.5), Inches(1.0), Inches(1.8), Inches(1.2), Inches(1.8), Inches(1.3)]
        for j, w in enumerate(widths):
            table.columns[j].width = w
    else:
        tb = slide2.shapes.add_textbox(Inches(3), Inches(3), Inches(7), Inches(1))
        p = tb.text_frame.paragraphs[0]
        p.text = "No trades today"
        p.font.size = Pt(24)
        p.font.color.rgb = MEDIUM_GRAY
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER

    # ─── Slide 3: Summary ─────────────────────────────────────────
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide3)

    # Logo
    circle = slide3.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.67), Inches(1.5), Inches(2), Inches(2))
    circle.fill.solid()
    circle.fill.fore_color.rgb = GOLD
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = "K"
    p.font.size = Pt(60)
    p.font.color.rgb = DARK_BG
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(8)

    tb = slide3.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11.33), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = f"Daily P/L: ${total_pnl:+,.2f} | {len(trades)} Trades | {len(wins)}W / {len(losses)}L"
    p.font.size = Pt(20)
    p.font.color.rgb = GREEN if total_pnl >= 0 else RED
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    tb2 = slide3.shapes.add_textbox(Inches(1), Inches(4.6), Inches(11.33), Inches(0.5))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = f"Balance: ${account.get('balance', 0):,.2f} | Equity: ${account.get('equity', 0):,.2f}"
    p2.font.size = Pt(14)
    p2.font.color.rgb = MEDIUM_GRAY
    p2.font.name = "Calibri"
    p2.alignment = PP_ALIGN.CENTER

    prs.save(output_path)
    return output_path


# ─── Generate & Send ──────────────────────────────────────────────
def generate_and_send_daily_report():
    """Generate daily PDF + PPTX and send to Telegram."""
    import logging
    log = logging.getLogger("daily_report")

    log.info("Generating daily report...")
    trades = get_today_trades()
    account = get_account_snapshot()

    now = datetime.now()
    date_tag = now.strftime("%Y%m%d")
    report_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "reports")
    os.makedirs(report_dir, exist_ok=True)

    # Generate PDF
    pdf_path = os.path.join(report_dir, f"Kingade_Daily_{date_tag}.pdf")
    generate_pdf(trades, account, pdf_path)
    log.info(f"PDF saved: {pdf_path}")

    # Generate PPTX
    pptx_path = os.path.join(report_dir, f"Kingade_Daily_{date_tag}.pptx")
    generate_pptx(trades, account, pptx_path)
    log.info(f"PPTX saved: {pptx_path}")

    # Send summary message
    total_pnl = sum(t["net_profit"] for t in trades)
    wins = [t for t in trades if t["net_profit"] >= 0]
    losses = [t for t in trades if t["net_profit"] < 0]

    summary = (
        f"<b>DAILY REPORT - {now.strftime('%d %b %Y')}</b>\n\n"
        f"<b>Trades:</b> {len(trades)}\n"
        f"<b>Wins:</b> {len(wins)} | <b>Losses:</b> {len(losses)}\n"
        f"<b>Daily P/L:</b> ${total_pnl:+,.2f}\n"
        f"<b>Balance:</b> ${account.get('balance', 0):,.2f}\n\n"
        f"<i>Sending reports below...</i>"
    )
    tg.send_message(summary)

    # Send files
    tg.send_document(pdf_path, caption=f"Kingade Daily Report PDF - {now.strftime('%d %b %Y')}")
    tg.send_document(pptx_path, caption=f"Kingade Daily Report PPTX - {now.strftime('%d %b %Y')}")

    log.info("Daily report sent to Telegram")
    return pdf_path, pptx_path


# ─── Weekly Report ─────────────────────────────────────────────────
def get_week_trades():
    """Get all trades from Monday to Friday of the current week.
    MT5 timestamps are UTC, so the week window must be UTC-based."""
    now = datetime.now(timezone.utc)
    monday = now - timedelta(days=now.weekday())
    monday_start = datetime(monday.year, monday.month, monday.day, tzinfo=timezone.utc)

    from_ts = int(monday_start.timestamp())
    to_ts = int(now.timestamp())

    deals = mt5.history_deals_get(from_ts, to_ts)
    if deals is None:
        return []

    bot_deals = [d for d in deals if d.magic == config.MAGIC_NUMBER]

    trades = []
    for d in bot_deals:
        trades.append({
            "ticket": d.order,
            "time": datetime.fromtimestamp(d.time),
            "symbol": d.symbol,
            "type": "BUY" if d.type == mt5.DEAL_TYPE_BUY else "SELL",
            "volume": d.volume,
            "price": d.price,
            "profit": d.profit,
            "swap": d.swap,
            "commission": d.commission,
            "net_profit": d.profit + d.swap + d.commission,
            "comment": d.comment,
        })

    return trades


def generate_weekly_pdf(trades, account, output_path):
    """Generate weekly PDF report."""
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle("Title2", parent=styles["Title"], fontSize=18, textColor=colors.HexColor("#1a1a2e")))
    styles.add(ParagraphStyle("Sub", parent=styles["Normal"], fontSize=10, textColor=colors.grey))

    doc = SimpleDocTemplate(output_path, pagesize=A4, leftMargin=15*mm, rightMargin=15*mm, topMargin=15*mm, bottomMargin=15*mm)
    elements = []

    now = datetime.now()
    week_num = now.isocalendar()[1]
    year = now.year
    monday = now - timedelta(days=now.weekday())
    friday = monday + timedelta(days=4)
    date_range = f"{monday.strftime('%d %b')} - {friday.strftime('%d %b %Y')}"

    elements.append(Paragraph("Kingade Scalper Bot", styles["Title2"]))
    elements.append(Paragraph(f"Weekly Report - Week {week_num} | {date_range}", styles["Sub"]))
    elements.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#16213e")))
    elements.append(Spacer(1, 8))

    wins = [t for t in trades if t["net_profit"] >= 0]
    losses = [t for t in trades if t["net_profit"] < 0]
    total_pnl = sum(t["net_profit"] for t in trades)
    gross_profit = sum(t["net_profit"] for t in wins) if wins else 0
    gross_loss = abs(sum(t["net_profit"] for t in losses)) if losses else 0
    pf = gross_profit / gross_loss if gross_loss > 0 else 0

    summary = [
        ["WEEKLY SUMMARY", "", "", ""],
        ["Week", f"Week {week_num}, {year}", "Date Range", date_range],
        ["Total Trades", str(len(trades)), "Win Rate", f"{len(wins)/len(trades)*100:.0f}%" if trades else "0%"],
        ["Wins", str(len(wins)), "Losses", str(len(losses))],
        ["Weekly P/L", f"${total_pnl:+,.2f}", "Profit Factor", f"{pf:.2f}"],
        ["Balance", f"${account.get('balance', 0):,.2f}", "Equity", f"${account.get('equity', 0):,.2f}"],
        ["Avg Win", f"${np.mean([t['net_profit'] for t in wins]):+,.2f}" if wins else "$0", "Avg Loss", f"${np.mean([t['net_profit'] for t in losses]):+,.2f}" if losses else "$0"],
    ]

    t = Table(summary, colWidths=[100, 140, 100, 140])
    t.setStyle(TableStyle([
        ("SPAN", (0, 0), (3, 0)),
        ("BACKGROUND", (0, 0), (3, 0), colors.HexColor("#16213e")),
        ("TEXTCOLOR", (0, 0), (3, 0), colors.white),
        ("FONTNAME", (0, 0), (3, 0), "Helvetica-Bold"),
        ("FONTSIZE", (0, 0), (3, 0), 11),
        ("ALIGN", (0, 0), (3, 0), "CENTER"),
        ("FONTNAME", (0, 1), (0, -1), "Helvetica-Bold"),
        ("FONTNAME", (2, 1), (2, -1), "Helvetica-Bold"),
        ("FONTSIZE", (0, 1), (-1, -1), 9),
        ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#cccccc")),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f5f5f5")]),
        ("TOPPADDING", (0, 0), (-1, -1), 4),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
    ]))
    elements.append(t)
    elements.append(Spacer(1, 12))

    # Day-by-day breakdown
    elements.append(Paragraph("Daily Breakdown", styles["Heading2"]))
    elements.append(HRFlowable(width="100%", thickness=0.5, color=colors.HexColor("#cccccc")))
    elements.append(Spacer(1, 4))

    day_names = ["Monday", "Tuesday", "Wednesday", "Thursday", "Friday"]
    day_header = ["Day", "Trades", "Wins", "Losses", "WR%", "P/L"]
    day_rows = [day_header]

    for i, day_name in enumerate(day_names):
        day_date = monday + timedelta(days=i)
        day_trades = [t for t in trades if t["time"].date() == day_date.date()]
        day_wins = [t for t in day_trades if t["net_profit"] >= 0]
        day_losses = [t for t in day_trades if t["net_profit"] < 0]
        day_pnl = sum(t["net_profit"] for t in day_trades)
        day_wr = f"{len(day_wins)/len(day_trades)*100:.0f}%" if day_trades else "-"

        day_rows.append([
            f"{day_name} {day_date.strftime('%d/%m')}",
            str(len(day_trades)),
            str(len(day_wins)),
            str(len(day_losses)),
            day_wr,
            f"${day_pnl:+,.2f}" if day_trades else "-",
        ])

    day_rows.append([
        "WEEKLY TOTAL",
        str(len(trades)),
        str(len(wins)),
        str(len(losses)),
        f"{len(wins)/len(trades)*100:.0f}%" if trades else "0%",
        f"${total_pnl:+,.2f}",
    ])

    t2 = Table(day_rows, colWidths=[110, 55, 50, 55, 50, 80])
    t2.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#0f3460")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("BACKGROUND", (0, -1), (-1, -1), colors.HexColor("#16213e")),
        ("TEXTCOLOR", (0, -1), (-1, -1), colors.white),
        ("FONTNAME", (0, -1), (-1, -1), "Helvetica-Bold"),
        ("FONTSIZE", (0, 0), (-1, -1), 8),
        ("ALIGN", (0, 0), (-1, -1), "CENTER"),
        ("GRID", (0, 0), (-1, -1), 0.3, colors.HexColor("#dddddd")),
        ("ROWBACKGROUNDS", (0, 1), (-1, -2), [colors.white, colors.HexColor("#f8f8f8")]),
        ("TOPPADDING", (0, 0), (-1, -1), 3),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
    ]))
    elements.append(t2)
    elements.append(Spacer(1, 12))

    # Per symbol breakdown
    elements.append(Paragraph("By Symbol", styles["Heading2"]))
    elements.append(HRFlowable(width="100%", thickness=0.5, color=colors.HexColor("#cccccc")))
    elements.append(Spacer(1, 4))

    sym_data = {}
    for t in trades:
        s = t["symbol"]
        if s not in sym_data:
            sym_data[s] = {"trades": 0, "wins": 0, "pnl": 0}
        sym_data[s]["trades"] += 1
        sym_data[s]["pnl"] += t["net_profit"]
        if t["net_profit"] >= 0:
            sym_data[s]["wins"] += 1

    sym_header = ["Symbol", "Trades", "Wins", "Losses", "WR%", "P/L"]
    sym_rows = [sym_header]
    for sym in sorted(sym_data.keys()):
        d = sym_data[sym]
        wr = f"{d['wins']/d['trades']*100:.0f}%" if d["trades"] else "0%"
        sym_rows.append([sym, str(d["trades"]), str(d["wins"]),
                        str(d["trades"] - d["wins"]), wr, f"${d['pnl']:+,.2f}"])

    t3 = Table(sym_rows, colWidths=[80, 55, 50, 55, 50, 80])
    t3.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#0f3460")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTSIZE", (0, 0), (-1, -1), 8),
        ("ALIGN", (0, 0), (-1, -1), "CENTER"),
        ("GRID", (0, 0), (-1, -1), 0.3, colors.HexColor("#dddddd")),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f8f8f8")]),
        ("TOPPADDING", (0, 0), (-1, -1), 3),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
    ]))
    elements.append(t3)

    doc.build(elements)
    return output_path


def generate_weekly_pptx(trades, account, output_path):
    """Generate weekly PPTX report."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    now = datetime.now()
    week_num = now.isocalendar()[1]
    monday = now - timedelta(days=now.weekday())
    friday = monday + timedelta(days=4)
    date_range = f"{monday.strftime('%d %b')} - {friday.strftime('%d %b %Y')}"

    wins = [t for t in trades if t["net_profit"] >= 0]
    losses = [t for t in trades if t["net_profit"] < 0]
    total_pnl = sum(t["net_profit"] for t in trades)
    gross_profit = sum(t["net_profit"] for t in wins) if wins else 0
    gross_loss = abs(sum(t["net_profit"] for t in losses)) if losses else 0
    pf = gross_profit / gross_loss if gross_loss > 0 else 0

    # Slide 1: Overview
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    _add_title_bar(slide, f"KINGADE SCALPER BOT - Weekly Report - Week {week_num}")

    y = Inches(1.2)
    _add_kpi(slide, Inches(0.5), y, "BALANCE", f"${account.get('balance', 0):,.2f}", WHITE, BLUE)
    _add_kpi(slide, Inches(2.9), y, "EQUITY", f"${account.get('equity', 0):,.2f}", WHITE, BLUE)
    _add_kpi(slide, Inches(5.3), y, "WEEKLY P/L", f"${total_pnl:+,.2f}", GREEN if total_pnl >= 0 else RED, GREEN if total_pnl >= 0 else RED)
    _add_kpi(slide, Inches(7.7), y, "TOTAL TRADES", str(len(trades)), WHITE, BLUE)
    _add_kpi(slide, Inches(10.1), y, "WIN RATE", f"{len(wins)/len(trades)*100:.0f}%" if trades else "0%", GREEN if trades and len(wins)/len(trades) >= 0.5 else RED, GREEN if trades and len(wins)/len(trades) >= 0.5 else RED)

    y2 = Inches(2.6)
    _add_kpi(slide, Inches(0.5), y2, "WINS", str(len(wins)), GREEN, GREEN)
    _add_kpi(slide, Inches(2.9), y2, "LOSSES", str(len(losses)), RED, RED)
    avg_win = np.mean([t["net_profit"] for t in wins]) if wins else 0
    avg_loss = np.mean([t["net_profit"] for t in losses]) if losses else 0
    _add_kpi(slide, Inches(5.3), y2, "AVG WIN", f"${avg_win:+,.2f}", GREEN, GREEN)
    _add_kpi(slide, Inches(7.7), y2, "AVG LOSS", f"${avg_loss:+,.2f}", RED, RED)
    _add_kpi(slide, Inches(10.1), y2, "PROFIT FACTOR", f"{pf:.2f}", GOLD, GOLD)

    # Slide 2: Daily Breakdown
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide2)
    _add_title_bar(slide2, f"DAILY BREAKDOWN - Week {week_num} | {date_range}")

    day_names = ["Monday", "Tuesday", "Wednesday", "Thursday", "Friday"]
    n_rows = 7
    tbl = slide2.shapes.add_table(n_rows, 6, Inches(0.5), Inches(1.2), Inches(12.3), Inches(4.0))
    table = tbl.table

    headers = ["Day", "Trades", "Wins", "Losses", "WR%", "P/L"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER

    for i, day_name in enumerate(day_names):
        day_date = monday + timedelta(days=i)
        day_trades = [t for t in trades if t["time"].date() == day_date.date()]
        day_wins_list = [t for t in day_trades if t["net_profit"] >= 0]
        day_losses_list = [t for t in day_trades if t["net_profit"] < 0]
        day_pnl = sum(t["net_profit"] for t in day_trades)
        day_wr = f"{len(day_wins_list)/len(day_trades)*100:.0f}%" if day_trades else "-"

        row_data = [
            f"{day_name} {day_date.strftime('%d/%m')}",
            str(len(day_trades)), str(len(day_wins_list)), str(len(day_losses_list)),
            day_wr, f"${day_pnl:+,.2f}" if day_trades else "-",
        ]
        for j, val in enumerate(row_data):
            cell = table.cell(i + 1, j)
            cell.text = val
            bg = RGBColor(0x1A, 0x1A, 0x2E) if i % 2 == 0 else DARK_BG
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.color.rgb = WHITE
                p.font.name = "Calibri"
                p.alignment = PP_ALIGN.CENTER
                if j == 5:
                    pnl_val = day_pnl if day_trades else 0
                    p.font.color.rgb = GREEN if pnl_val >= 0 else RED
                    p.font.bold = True

    # Total row
    total_data = ["WEEKLY TOTAL", str(len(trades)), str(len(wins)), str(len(losses)),
                   f"{len(wins)/len(trades)*100:.0f}%" if trades else "0%",
                   f"${total_pnl:+,.2f}"]
    for j, val in enumerate(total_data):
        cell = table.cell(6, j)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.color.rgb = GOLD
            p.font.bold = True
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER

    # Slide 3: By Symbol
    sym_data = {}
    for t in trades:
        s = t["symbol"]
        if s not in sym_data:
            sym_data[s] = {"trades": 0, "wins": 0, "pnl": 0}
        sym_data[s]["trades"] += 1
        sym_data[s]["pnl"] += t["net_profit"]
        if t["net_profit"] >= 0:
            sym_data[s]["wins"] += 1

    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide3)
    _add_title_bar(slide3, "BY SYMBOL")

    sym_list = sorted(sym_data.keys())
    n_sym = len(sym_list) + 1
    tbl3 = slide3.shapes.add_table(n_sym, 6, Inches(0.5), Inches(1.2), Inches(12.3), Inches(4.0))
    table3 = tbl3.table

    sym_headers = ["Symbol", "Trades", "Wins", "Losses", "WR%", "P/L"]
    for j, h in enumerate(sym_headers):
        cell = table3.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER

    for i, sym in enumerate(sym_list):
        d = sym_data[sym]
        wr = f"{d['wins']/d['trades']*100:.0f}%" if d["trades"] else "0%"
        row_data = [sym, str(d["trades"]), str(d["wins"]),
                   str(d["trades"] - d["wins"]), wr, f"${d['pnl']:+,.2f}"]
        for j, val in enumerate(row_data):
            cell = table3.cell(i + 1, j)
            cell.text = val
            bg = RGBColor(0x1A, 0x1A, 0x2E) if i % 2 == 0 else DARK_BG
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.color.rgb = WHITE
                p.font.name = "Calibri"
                p.alignment = PP_ALIGN.CENTER
                if j == 5:
                    p.font.color.rgb = GREEN if d["pnl"] >= 0 else RED
                    p.font.bold = True

    prs.save(output_path)
    return output_path


# ─── Generate & Send Weekly ───────────────────────────────────────
def generate_and_send_weekly_report():
    """Generate weekly PDF + PPTX and send to Telegram."""
    import logging
    log = logging.getLogger("daily_report")

    log.info("Generating weekly report...")
    trades = get_week_trades()
    account = get_account_snapshot()

    now = datetime.now()
    week_num = now.isocalendar()[1]
    year = now.year
    report_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "reports")
    os.makedirs(report_dir, exist_ok=True)

    # Generate PDF
    pdf_path = os.path.join(report_dir, f"Kingade_Weekly_{year}-W{week_num:02d}.pdf")
    generate_weekly_pdf(trades, account, pdf_path)
    log.info(f"Weekly PDF saved: {pdf_path}")

    # Generate PPTX
    pptx_path = os.path.join(report_dir, f"Kingade_Weekly_{year}-W{week_num:02d}.pptx")
    generate_weekly_pptx(trades, account, pptx_path)
    log.info(f"Weekly PPTX saved: {pptx_path}")

    # Send summary message
    total_pnl = sum(t["net_profit"] for t in trades)
    wins = [t for t in trades if t["net_profit"] >= 0]
    losses = [t for t in trades if t["net_profit"] < 0]

    summary = (
        f"<b>WEEKLY REPORT - Week {week_num}, {year}</b>\n\n"
        f"<b>Trades:</b> {len(trades)}\n"
        f"<b>Wins:</b> {len(wins)} | <b>Losses:</b> {len(losses)}\n"
        f"<b>Weekly P/L:</b> ${total_pnl:+,.2f}\n"
        f"<b>Balance:</b> ${account.get('balance', 0):,.2f}\n\n"
        f"<i>Sending weekly reports below...</i>"
    )
    tg.send_message(summary)

    # Send files
    tg.send_document(pdf_path, caption=f"Kingade Weekly Report PDF - Week {week_num}")
    tg.send_document(pptx_path, caption=f"Kingade Weekly Report PPTX - Week {week_num}")

    log.info("Weekly report sent to Telegram")
    return pdf_path, pptx_path

